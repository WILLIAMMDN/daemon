import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape

images = [
    r"C:\Users\MEDINA\Downloads\ChatGPT Image 27 jul 2026, 01_40_51 a.m. (1).png",
    r"C:\Users\MEDINA\Downloads\ChatGPT Image 27 jul 2026, 01_40_51 a.m. (2).png",
    r"C:\Users\MEDINA\Downloads\ChatGPT Image 27 jul 2026, 01_40_51 a.m. (3).png",
    r"C:\Users\MEDINA\Downloads\ChatGPT Image 27 jul 2026, 01_40_52 a.m. (4).png",
    r"C:\Users\MEDINA\Downloads\ChatGPT Image 27 jul 2026, 01_40_52 a.m. (5).png"
]

output_pptx = r"C:\Users\MEDINA\Desktop\daemon-aula-sustentacion\scratch\Letraria_Presentacion.pptx"
output_pdf = r"C:\Users\MEDINA\Desktop\daemon-aula-sustentacion\scratch\Letraria_Presentacion.pdf"

# --- PPTX Generation ---
prs = Presentation()
# Set presentation slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
for img_path in images:
    if not os.path.exists(img_path):
        print(f"Not found: {img_path}")
        continue
    slide = prs.slides.add_slide(blank_slide_layout)
    # Add image covering the entire slide
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(output_pptx)
print("Saved PPTX.")

# --- PDF Generation ---
# The images are 16:9, so we'll use a 16:9 page size in points (1 point = 1/72 inch)
width, height = 13.333 * 72, 7.5 * 72 
c = canvas.Canvas(output_pdf, pagesize=(width, height))

for img_path in images:
    if not os.path.exists(img_path):
        continue
    # Draw image filling the page
    c.drawImage(img_path, 0, 0, width=width, height=height)
    c.showPage()

c.save()
print("Saved PDF.")
